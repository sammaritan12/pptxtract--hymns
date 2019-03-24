from glob import glob
from sys import argv
from re import sub
from os.path import basename

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor


def list_filenames(dirpath, file_type):
    if type(dirpath) is not str or type(file_type) is not str:
        raise TypeError('dirpath and filetype must both be strings')

    return glob(dirpath + '/**/*.' + file_type, recursive=True)


def filenames_only(filenames):
    named_files = []
    for i in filenames:
        for j in range(len(i)):
            if i[j] == '/':
                named_files.append(i[j+1:])
                continue
            elif i[j] == '\\':
                named_files.append(i[j+1:])
                continue
    return named_files

if len(argv) < 6:
    print("Enter the following args")
    print("extract.py <dir to extract> <dir to place> <background images> <text colour> <wide/normal>")
    quit()

# Argument variables
extract_dir = argv[1]
save_dir = argv[2]
image_dir = argv[3]
text_colour = argv[4]
aspect_ratio = argv[5]

# Get all pptx files and image files
pptx_file_dir = list_filenames(extract_dir, 'pptx')
# pptx_file_names = filenames_only(pptx_file_dir)
pptx_file_names =  list(map(lambda x: basename(x), pptx_file_dir))
images_path = list_filenames(image_dir, 'jpg')

# List of hymns
hymns = []

# Change colour, normally black if none supplied
text_colour = RGBColor(255, 255, 255) if text_colour.lower() == 'white' else RGBColor(0,0,0)

# Extract the hymn texts into a list of strings
for i in range(len(pptx_file_dir)):
    current_hymn = []
    prs = Presentation(pptx_file_dir[i])

    # Choose only objects which have text
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text.strip() == '':
                # current_hymn.append(shape.text)
                current_hymn.append(sub(' +', ' ', '\n'.join(list(map(lambda x: x.strip(), shape.text.split('\n'))))).strip(" \n\t*"))
    hymns.append(current_hymn)

# See if any extracfted hymns are empty
for j in range(len(hymns)):
    if len(hymns[j]) == 0:
        print(pptx_file_names[j], "Length: 0")

# Start creating powerpoints for all of the extracted hymns
for h in range(len(hymns)):

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    if aspect_ratio.lower() == 'wide':
        # 16:9 Aspect ratio
        prs.slide_height = 6858000
        prs.slide_width = 12192000
    else:
        # 4:3 Aspect ration
        prs.slide_height = 6858000
        prs.slide_width = 9144000

    # Create Title Slide
    title_slide = prs.slides.add_slide(blank_slide_layout)

    # Add picture
    pic = title_slide.shapes.add_picture(images_path[h % len(images_path)], 0, 0)
    pic.height = prs.slide_height
    pic.width = prs.slide_width

    # Add title
    left = Mm(20)
    top = 2130425
    width = prs.slide_width - Mm(40)
    height = Mm(60)

    title = title_slide.shapes.add_textbox(left, top,  width , height)
    title.text = hymns[h][0]

    title.text_frame.word_wrap = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.color.rgb = text_colour

    for i in range(1, len(hymns[h])):
        # Create Blank slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add background image
        pic = slide.shapes.add_picture(images_path[h % len(images_path)], 0, 0, width=prs.slide_width , height=prs.slide_height)

        # 5mm from top and left
        left = top = Mm(5)

        # Encompass width and height of all
        width = prs.slide_width - Mm(10)
        height = prs.slide_height- Mm(10)
        
        # Add lyrics
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = hymns[h][i] + ' *' if i == len(hymns[h])-1 else hymns[h][i]

        # Change size, font and colour
        tf.paragraphs[0].font.name = 'Arial'
        tf.paragraphs[0].font.size = Pt(34)
        tf.paragraphs[0].font.color.rgb = text_colour
        tf.paragraphs[0].font.bold = False
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT   

    print('{}/{} {} Processed'.format(h+1, len(hymns), pptx_file_names[h]))

    # Save file
    prs.save(save_dir+'/'+pptx_file_names[h])
